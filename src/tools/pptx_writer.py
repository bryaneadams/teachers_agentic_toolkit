# type: ignore
from __future__ import annotations

import re
from pathlib import Path

from errors import RenderingError
from models import Slide, SlideDeck


def write_pptx(
    deck: SlideDeck,
    output_path: str | Path,
    template_path: str | Path | None = None,
) -> tuple[Path, tuple[str, ...]]:
    """Saves your PowerPoint

    Args:
        deck (SlideDeck): slide deck to save
        output_path (str | Path): output directory
        template_path (str | Path | None, optional): Path to template. Defaults to None.

    Raises:
        RenderingError: Could not write PowerPoint file
        RenderingError: Could not render PowerPoint file

    Returns:
        tuple[Path, tuple[str, ...]]: save path and warnings
    """
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    warnings: list[str] = []

    try:
        presentation = _presentation(template_path)
        _clear_default_slides(presentation)
        for slide in deck.slides:
            _add_content_slide(presentation, slide, warnings)
        presentation.save(output)
    except OSError as exc:
        raise RenderingError(f"Could not write PowerPoint file: {output}") from exc
    except Exception as exc:
        raise RenderingError(f"Could not render PowerPoint file: {exc}") from exc

    return output, tuple(warnings)


def _presentation(template_path: str | Path | None):
    """Creates the presentation format

    Args:
        template_path (str | Path | None): path to template to use

    Raises:
        RenderingError: could not find template path

    Returns:
        Presentation: returns the presentation object
    """
    from pptx import Presentation
    from pptx.util import Inches

    if template_path is not None:
        template = Path(template_path)
        if not template.exists():
            raise RenderingError(f"Template does not exist: {template}")
        return Presentation(str(template))

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    return presentation


def _clear_default_slides(presentation) -> None:
    """Removes default slides

    Args:
        presentation (Presentation): presentation object
    """
    slide_ids = list(
        presentation.slides._sldIdLst
    )  # noqa: SLF001 - python-pptx has no public delete API.
    for slide_id in slide_ids:
        rel_id = slide_id.rId
        presentation.part.drop_rel(rel_id)
        presentation.slides._sldIdLst.remove(slide_id)  # noqa: SLF001


def _add_content_slide(presentation, slide: Slide, warnings: list[str]) -> None:
    """Adds content to a slide in the correct format

    Args:
        presentation (Presentation): Presentation opbject
        slide (Slide): slide to update
        warnings (list[str]): any warnings
    """
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    layout = _blank_layout(presentation)
    pptx_slide = presentation.slides.add_slide(layout)

    title_box = pptx_slide.shapes.add_textbox(
        Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.clear()
    title_para = title_frame.paragraphs[0]
    title_para.text = _plain_slide_text(slide.title)
    title_para.alignment = PP_ALIGN.LEFT
    title_run = title_para.runs[0]
    title_run.font.bold = True
    title_run.font.size = Pt(30)

    has_image = bool(slide.image_path and slide.image_path.exists())
    body_width = Inches(7.1 if has_image else 12.0)
    body_box = pptx_slide.shapes.add_textbox(
        Inches(0.75), Inches(1.4), body_width, Inches(4.8)
    )
    body_frame = body_box.text_frame
    body_frame.clear()
    body_frame.word_wrap = True

    for index, bullet in enumerate(slide.bullets):
        paragraph = (
            body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
        )
        paragraph.text = _plain_slide_text(bullet)
        paragraph.level = 0
        _enable_bullet(paragraph)
        paragraph.font.size = Pt(20)

    if has_image and slide.image_path is not None:
        _add_image(pptx_slide, Path(slide.image_path), warnings)
    elif slide.image_path:
        warnings.append(f"Image was not found and was skipped: {slide.image_path}")

    if slide.speaker_notes:
        _add_speaker_notes(pptx_slide, _plain_slide_text(slide.speaker_notes))


def _blank_layout(presentation):
    """Finds the blank layout

    Args:
        presentation (Presentation): Presentation object

    Returns:
        Presentation: blank slide layout
    """
    for layout in presentation.slide_layouts:
        if layout.name.lower() == "blank":
            return layout
    return presentation.slide_layouts[-1]


def _add_image(pptx_slide, image_path: Path, warnings: list[str]) -> None:
    """Adds images to the slide

    Args:
        pptx_slide (Presentation): presentation object
        image_path (Path): path to image
        warnings (list[str]): any warnings
    """
    from pptx.util import Inches

    try:
        pptx_slide.shapes.add_picture(
            str(image_path),
            Inches(8.25),
            Inches(1.55),
            width=Inches(4.45),
            height=Inches(4.45),
        )
    except Exception as exc:
        warnings.append(
            f"Image could not be added and was skipped: {image_path} ({exc})"
        )


def _add_speaker_notes(pptx_slide, notes: str) -> None:
    """Adds speaker notes to bottom of the slide

    Args:
        pptx_slide (Presentation): presentation object
        notes (str): notes to add to the slide for speaker
    """
    notes_frame = pptx_slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = notes


def _enable_bullet(paragraph) -> None:
    """Adds bullets to slide

    Args:
        paragraph (Presentation): Presentation object
    """
    from pptx.oxml.xmlchemy import OxmlElement

    properties = paragraph._p.get_or_add_pPr()  # noqa: SLF001 - no public bullet API.
    for existing in properties.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar"
    ):
        properties.remove(existing)
    bullet = OxmlElement("a:buChar")
    bullet.set("char", "\u2022")
    properties.insert(0, bullet)


def _plain_slide_text(text: str) -> str:
    """Removes unnecessary text from slide content

    Args:
        text (str): slide text

    Returns:
        str: cleaned text
    """
    clean = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", text)
    clean = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", clean)
    clean = re.sub(r"[*_`~]+", "", clean)
    clean = re.sub(r"^\s*(?:[-*+]|\d+[.)]|\u2022)\s+", "", clean)
    clean = re.sub(r"\s+", " ", clean)
    return clean.strip()
